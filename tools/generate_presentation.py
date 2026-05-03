# -*- coding: utf-8 -*-
"""Sinh file PowerPoint thuyết trình cho đồ án VocabQuest.

Chạy:
    python tools/generate_presentation.py

Output:
    C:/Users/Kieu Anh/Desktop/CD2/VocabQuest_Presentation.pptx
"""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# -------------------- BẢNG MÀU & FONT --------------------
PRIMARY = RGBColor(0x6C, 0x5C, 0xE7)        # tím
PRIMARY_DARK = RGBColor(0x4C, 0x3F, 0xB8)
ACCENT = RGBColor(0xFF, 0x6B, 0x6B)         # cam san hô
ACCENT2 = RGBColor(0x00, 0xC6, 0xA7)        # ngọc lục
GOLD = RGBColor(0xFF, 0xC1, 0x07)
BG_LIGHT = RGBColor(0xF6, 0xF7, 0xFB)
BG_DARK = RGBColor(0x1F, 0x1B, 0x2E)
TEXT_DARK = RGBColor(0x1B, 0x1D, 0x2A)
TEXT_MUTED = RGBColor(0x6B, 0x70, 0x83)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x10, 0x10, 0x18)

FONT_TITLE = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# -------------------- HELPER --------------------
def add_rect(slide, x, y, w, h, fill, line=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape


def add_round_rect(slide, x, y, w, h, fill, corner=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = corner
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def set_solid_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def set_gradient_bg(slide, color1, color2, angle=45):
    """Background gradient bằng cách insert XML trực tiếp."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    sppr = bg.fill._xPr
    for child in list(sppr):
        if child.tag.endswith('}solidFill') or child.tag.endswith('}gradFill') \
                or child.tag.endswith('}noFill') or child.tag.endswith('}blipFill') \
                or child.tag.endswith('}pattFill'):
            sppr.remove(child)
    grad = etree.SubElement(sppr, qn('a:gradFill'))
    grad.set('flip', 'none')
    grad.set('rotWithShape', '1')
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, col in ((0, color1), (100000, color2)):
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(pos))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(col[0], col[1], col[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(angle * 60000))
    lin.set('scaled', '1')
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=18, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT_DARK,
                bullet_color=PRIMARY, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r1 = p.add_run()
        r1.text = "●  "
        r1.font.name = FONT_BODY
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT_BODY
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_title_bar(slide, title, subtitle=None, *, dark=False):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.95))
    bar.line.fill.background()
    sppr = bar.fill._xPr
    for child in list(sppr):
        if child.tag.endswith('}solidFill') or child.tag.endswith('}gradFill'):
            sppr.remove(child)
    grad = etree.SubElement(sppr, qn('a:gradFill'))
    grad.set('flip', 'none')
    grad.set('rotWithShape', '1')
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    c1 = (0x6C, 0x5C, 0xE7)
    c2 = (0xFF, 0x6B, 0x6B)
    for pos, col in ((0, c1), (100000, c2)):
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(pos))
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(*col))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(0))
    lin.set('scaled', '1')

    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.4), Inches(0.27), Inches(0.4), Inches(0.4))
    dot.fill.solid()
    dot.fill.fore_color.rgb = WHITE
    dot.line.fill.background()
    dot2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.55), Inches(0.42), Inches(0.16), Inches(0.16))
    dot2.fill.solid()
    dot2.fill.fore_color.rgb = ACCENT
    dot2.line.fill.background()

    add_text(slide, Inches(1.0), Inches(0.18), Inches(11.5), Inches(0.55),
             title, font=FONT_TITLE, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(1.0), Inches(0.62), Inches(11.5), Inches(0.4),
                 subtitle, font=FONT_BODY, size=13, color=WHITE, italic=True)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.95), SLIDE_W, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()


def add_footer(slide, page_no, total):
    add_text(slide, Inches(0.4), Inches(7.15), Inches(6.0), Inches(0.3),
             "VocabQuest | Báo cáo cuối môn", size=11, color=TEXT_MUTED,
             italic=True)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             f"{page_no} / {total}", size=11, color=TEXT_MUTED,
             align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, *, color=PRIMARY, icon=None):
    """Card có bo trái màu, tiêu đề, mô tả."""
    card = add_round_rect(slide, x, y, w, h, WHITE, corner=0.08)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.12), h)
    bar.adjustments[0] = 0.5
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    if icon:
        ic = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.25),
            Inches(0.6), Inches(0.6))
        ic.fill.solid()
        ic.fill.fore_color.rgb = color
        ic.line.fill.background()
        tf = ic.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = icon
        r.font.name = FONT_TITLE
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = WHITE
        title_x = x + Inches(1.0)
        title_w = w - Inches(1.2)
    else:
        title_x = x + Inches(0.3)
        title_w = w - Inches(0.4)
    add_text(slide, title_x, y + Inches(0.22), title_w, Inches(0.4),
             title, font=FONT_TITLE, size=15, bold=True, color=color)
    add_text(slide, x + Inches(0.3), y + Inches(0.85),
             w - Inches(0.5), h - Inches(1.0),
             body, size=12, color=TEXT_DARK)


def add_step_chip(slide, x, y, num, text, color):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(num)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.size = Pt(18)
    r.font.name = FONT_TITLE
    add_text(slide, x + Inches(0.7), y + Inches(0.05),
             Inches(4.5), Inches(0.5), text,
             size=14, bold=True, color=TEXT_DARK)


# -------------------- TỪNG SLIDE --------------------
def slide_cover(prs, total_slides, logo_path, info):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, (0x4C, 0x3F, 0xB8), (0xFF, 0x6B, 0x6B), angle=120)

    for cx, cy, cw, alpha in [
        (Inches(11.0), Inches(-1.5), Inches(5.0), None),
        (Inches(-1.0), Inches(5.5), Inches(4.0), None),
        (Inches(9.5), Inches(5.5), Inches(2.5), None),
    ]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, cw, cw)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.fill.background()
        sppr = c.fill._xPr
        solid = sppr.find(qn('a:solidFill'))
        if solid is not None:
            srgb = solid.find(qn('a:srgbClr'))
            if srgb is not None:
                a = etree.SubElement(srgb, qn('a:alpha'))
                a.set('val', '12000')

    if logo_path and Path(logo_path).exists():
        try:
            s.shapes.add_picture(str(logo_path),
                                 Inches(0.7), Inches(0.6),
                                 width=Inches(1.5))
        except Exception:
            pass

    add_text(s, Inches(0.7), Inches(2.4), Inches(11.0), Inches(0.5),
             "BÁO CÁO ĐỒ ÁN CUỐI MÔN", font=FONT_BODY, size=16,
             bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(2.95), Inches(11.5), Inches(1.5),
             "VocabQuest", font=FONT_TITLE, size=72, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(4.2), Inches(11.5), Inches(0.7),
             "Ứng dụng học từ vựng tiếng Anh qua mini-game",
             font=FONT_TITLE, size=24, color=WHITE)
    add_text(s, Inches(0.7), Inches(4.85), Inches(11.5), Inches(0.5),
             "Flutter  -  Dart  -  Lập trình bất đồng bộ  -  Firebase",
             font=FONT_BODY, size=15, color=WHITE, italic=True)

    box = add_round_rect(s, Inches(0.7), Inches(5.7),
                         Inches(11.9), Inches(1.3), WHITE, corner=0.12)
    sppr = box.fill._xPr
    solid = sppr.find(qn('a:solidFill'))
    if solid is not None:
        srgb = solid.find(qn('a:srgbClr'))
        if srgb is not None:
            a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', '20000')
    add_text(s, Inches(1.0), Inches(5.85), Inches(11.5), Inches(0.4),
             info.get("subject", ""), size=14, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(6.25), Inches(11.5), Inches(0.4),
             "Sinh viên thực hiện: " + info.get("student", ""),
             size=13, color=WHITE)
    add_text(s, Inches(1.0), Inches(6.55), Inches(11.5), Inches(0.4),
             "Giảng viên: " + info.get("teacher", "") +
             "    |    " + info.get("year", ""),
             size=13, color=WHITE)


def slide_toc(prs, total, sections):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "Mục lục", "Cấu trúc bài thuyết trình")

    col_w = Inches(5.8)
    row_h = Inches(0.55)
    half = (len(sections) + 1) // 2
    for i, item in enumerate(sections):
        col = 0 if i < half else 1
        row = i if i < half else i - half
        x = Inches(0.6) + col * (col_w + Inches(0.4))
        y = Inches(1.4) + row * row_h

        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4))
        c.fill.solid()
        c.fill.fore_color.rgb = PRIMARY if col == 0 else ACCENT
        c.line.fill.background()
        tf = c.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.size = Pt(13)
        r.font.name = FONT_TITLE

        add_text(s, x + Inches(0.55), y + Inches(0.04),
                 col_w - Inches(0.6), Inches(0.4),
                 item, size=14, bold=True, color=TEXT_DARK)
    add_footer(s, 2, total)


def slide_intro(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "1. Giới thiệu đề tài",
                  "Bài toán, lý do chọn và đối tượng người dùng")

    add_text(s, Inches(0.6), Inches(1.3), Inches(7.5), Inches(0.5),
             "Bài toán", font=FONT_TITLE, size=20, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.6), Inches(1.85), Inches(7.5), Inches(2.8), [
        "Học từ vựng tiếng Anh truyền thống dễ gây nhàm chán, khó duy trì",
        "Người học thiếu động lực và công cụ theo dõi tiến độ",
        "Cần một ứng dụng di động: vừa học, vừa chơi, vừa thưởng",
        "Yêu cầu: tương tác thời gian thực, hoạt động on/offline",
    ], size=14)

    add_text(s, Inches(0.6), Inches(4.7), Inches(7.5), Inches(0.5),
             "Mục tiêu của đề tài", font=FONT_TITLE, size=20,
             bold=True, color=ACCENT)
    add_bullets(s, Inches(0.6), Inches(5.25), Inches(7.5), Inches(1.8), [
        "Xây dựng app Flutter đa nền tảng (Android - iOS - Web)",
        "4 mini-game: Matching, Quiz, Word Puzzle, Memory",
        "Tích hợp Firebase (Auth, Firestore, Storage)",
        "Hệ thống XP - Level - Coin - Streak - Bảng xếp hạng",
    ], size=14)

    add_card(s, Inches(8.6), Inches(1.3), Inches(4.2), Inches(1.6),
             "Đối tượng sử dụng",
             "Học sinh - sinh viên và người đi làm từ trình độ A1 - B2.\n"
             "Muốn học từ vựng qua game ngắn, vui vẻ.",
             color=ACCENT2, icon="U")
    add_card(s, Inches(8.6), Inches(3.05), Inches(4.2), Inches(1.6),
             "Phạm vi",
             "3 gói từ (Beginner / Intermediate / Advanced).\n"
             "4 mini-game x 3 mức khó mỗi level.",
             color=PRIMARY, icon="P")
    add_card(s, Inches(8.6), Inches(4.8), Inches(4.2), Inches(1.6),
             "Công nghệ lõi",
             "Flutter 3.x - Dart 3 - Provider - Firebase.\n"
             "Lập trình bất đồng bộ với Future / Stream / async.",
             color=GOLD, icon="T")
    add_footer(s, 3, total)


def slide_flutter_theory(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "2. Lý thuyết - Flutter",
                  "Framework UI đa nền tảng của Google")

    add_text(s, Inches(0.6), Inches(1.25), Inches(7.5), Inches(0.5),
             "Flutter là gì?", font=FONT_TITLE, size=22, bold=True,
             color=PRIMARY)
    add_text(s, Inches(0.6), Inches(1.8), Inches(7.5), Inches(2.0),
             "Flutter là SDK mã nguồn mở, dùng ngôn ngữ Dart, cho phép "
             "xây dựng ứng dụng iOS - Android - Web - Desktop từ một bộ "
             "mã nguồn duy nhất. Flutter vẽ UI bằng Skia/Impeller, không "
             "phụ thuộc widget của hệ điều hành nên giao diện thống nhất "
             "trên mọi nền tảng.",
             size=14, color=TEXT_DARK)

    add_text(s, Inches(0.6), Inches(3.75), Inches(7.5), Inches(0.5),
             "Đặc điểm chính", font=FONT_TITLE, size=22, bold=True,
             color=ACCENT)
    add_bullets(s, Inches(0.6), Inches(4.3), Inches(7.5), Inches(2.6), [
        "Tất cả là Widget: StatelessWidget vs StatefulWidget",
        "Khai báo UI bằng Dart (declarative) - không dùng XML",
        "Hot Reload / Hot Restart, vòng đời phát triển siêu nhanh",
        "Render engine riêng - không phụ thuộc native widget",
        "Material Design + Cupertino sẵn có",
    ], size=14)

    box = add_round_rect(s, Inches(8.6), Inches(1.25),
                         Inches(4.2), Inches(5.7), WHITE, corner=0.06)
    add_text(s, Inches(8.85), Inches(1.4), Inches(3.9), Inches(0.4),
             "Cây Widget mẫu", size=14, bold=True, color=PRIMARY)
    tree = [
        ("MaterialApp", PRIMARY, 0),
        ("Scaffold", PRIMARY_DARK, 1),
        ("AppBar", ACCENT, 2),
        ("Body: Column", ACCENT, 2),
        ("Hero (Avatar)", ACCENT2, 3),
        ("StatsStrip", ACCENT2, 3),
        ("MainActions", ACCENT2, 3),
        ("Game Card x N", GOLD, 4),
        ("BottomNav", PRIMARY_DARK, 1),
    ]
    for i, (lbl, col, lvl) in enumerate(tree):
        y = Inches(1.85) + i * Inches(0.5)
        x = Inches(8.85) + lvl * Inches(0.25)
        chip = add_round_rect(s, x, y,
                              Inches(3.5) - lvl * Inches(0.2),
                              Inches(0.4), col, corner=0.4)
        tf = chip.text_frame
        tf.margin_left = Pt(8)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = lbl
        r.font.name = FONT_BODY
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE
    add_footer(s, 4, total)


def slide_dart_theory(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "3. Lý thuyết - Dart",
                  "Ngôn ngữ nền tảng của Flutter")

    add_text(s, Inches(0.6), Inches(1.25), Inches(7.5), Inches(0.5),
             "Dart - đặc trưng", font=FONT_TITLE, size=22, bold=True,
             color=PRIMARY)
    add_bullets(s, Inches(0.6), Inches(1.85), Inches(7.5), Inches(3.5), [
        "Ngôn ngữ hiện đại, type-safe, null-safety mặc định (Dart 3)",
        "Biên dịch AOT cho mobile, JIT khi dev (Hot Reload)",
        "Single-thread + Event Loop, hỗ trợ Isolate khi cần song song",
        "Cú pháp quen thuộc kiểu Java/JS - dễ tiếp cận",
        "Có tính năng: Future, Stream, async/await, mixin, extension",
        "Quản lý bộ nhớ qua garbage collector",
    ], size=14)

    code_box = add_round_rect(s, Inches(8.5), Inches(1.25),
                              Inches(4.3), Inches(5.7), INK, corner=0.04)
    add_text(s, Inches(8.7), Inches(1.4), Inches(4.0), Inches(0.4),
             "Ví dụ Dart", size=12, bold=True, color=GOLD)
    code = (
        "class UserModel {\n"
        "  final String uid;\n"
        "  final int totalXP;\n"
        "  final int streak;\n\n"
        "  UserModel({\n"
        "    required this.uid,\n"
        "    this.totalXP = 0,\n"
        "    this.streak = 0,\n"
        "  });\n\n"
        "  int get level =>\n"
        "    LevelSystem.fromXP(totalXP);\n"
        "}\n\n"
        "Future<UserModel?> getUser(String id) async {\n"
        "  final doc = await db.doc(id).get();\n"
        "  return UserModel.fromMap(doc.data());\n"
        "}"
    )
    tb = s.shapes.add_textbox(Inches(8.7), Inches(1.85),
                              Inches(4.0), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.name = FONT_MONO
    r.font.size = Pt(11)
    r.font.color.rgb = WHITE
    add_footer(s, 5, total)


def slide_async_theory(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "4. Lập trình bất đồng bộ trong Dart",
                  "Future - async/await - Stream")

    add_text(s, Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.5),
             "Vì sao cần bất đồng bộ?", size=18, bold=True, color=PRIMARY)
    add_text(s, Inches(0.6), Inches(1.7), Inches(12.0), Inches(1.2),
             "Các tác vụ như gọi mạng, đọc file, query Firestore có thể "
             "tốn hàng giây. Nếu chạy tuần tự sẽ khóa luồng UI -> màn hình "
             "đờ băng. Bất đồng bộ cho phép tác vụ chạy nền, UI vẫn phản "
             "hồi mượt mà.", size=13, color=TEXT_DARK)

    base_y = Inches(2.95)
    col_w = Inches(4.05)
    gap = Inches(0.15)
    x0 = Inches(0.55)

    add_card(s, x0, base_y, col_w, Inches(3.7),
             "Future<T>",
             "Đại diện cho 1 giá trị sẽ có ở tương lai. Dùng khi cần "
             "kết quả một lần: đăng nhập, đọc Firestore một lần, ghi file.\n"
             "API: .then() / await / FutureBuilder.\n"
             "Ví dụ: UserModel? u = await getUser(uid);",
             color=PRIMARY, icon="F")
    add_card(s, x0 + col_w + gap, base_y, col_w, Inches(3.7),
             "async / await",
             "Từ khóa giúp viết code bất đồng bộ theo phong cách tuần tự, "
             "dễ đọc, dễ bắt lỗi.\n"
             "Hàm có async trả về Future. await tạm dừng đến khi có kết quả "
             "mà không khóa thread chính.",
             color=ACCENT, icon="A")
    add_card(s, x0 + 2 * (col_w + gap), base_y, col_w, Inches(3.7),
             "Stream<T>",
             "Đẩy các giá trị liên tiếp theo thời gian. Dùng khi cần "
             "lắng nghe liên tục: Firestore realtime, sự kiện input, audio.\n"
             "API: .listen() / await for / StreamBuilder.\n"
             "VocabQuest dùng cho user/leaderboard realtime.",
             color=ACCENT2, icon="S")

    add_footer(s, 6, total)


def slide_event_loop(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "5. Event Loop & Isolate",
                  "Cách Dart xử lý bất đồng bộ")

    add_text(s, Inches(0.6), Inches(1.2), Inches(7.0), Inches(0.5),
             "Event Loop trong Dart", size=20, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.6), Inches(1.75), Inches(7.0), Inches(3.0), [
        "Mỗi Isolate có 1 thread + 1 Event Loop riêng",
        "Microtask Queue (ưu tiên cao) -> Event Queue",
        "Khi await xong, callback được đẩy vào queue để chạy tiếp",
        "Nhờ đó UI không bị giật khi có I/O",
    ], size=14)

    add_text(s, Inches(0.6), Inches(4.85), Inches(7.0), Inches(0.5),
             "Isolate - khi nào cần?", size=20, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.6), Inches(5.4), Inches(7.0), Inches(1.7), [
        "Tính toán nặng: parse JSON lớn, mã hóa, xử lý ảnh",
        "Dùng compute() để tách sang isolate phụ",
        "VocabQuest hiện dùng event loop là đủ, JSON ~3 MB",
    ], size=14)

    box = add_round_rect(s, Inches(7.9), Inches(1.2),
                         Inches(5.0), Inches(5.7), WHITE, corner=0.05)
    add_text(s, Inches(8.1), Inches(1.35), Inches(4.6), Inches(0.4),
             "Mô hình Event Loop", size=14, bold=True, color=PRIMARY)
    items = [
        ("UI - Render Frame", PRIMARY),
        ("Microtask Queue", ACCENT2),
        ("Event Queue (await, IO)", ACCENT),
        ("Future complete", GOLD),
        ("setState / notifyListeners", PRIMARY_DARK),
        ("Repaint Widget Tree", ACCENT),
    ]
    for i, (lbl, col) in enumerate(items):
        y = Inches(1.95) + i * Inches(0.78)
        chip = add_round_rect(s, Inches(8.1), y,
                              Inches(4.6), Inches(0.55), col, corner=0.4)
        tf = chip.text_frame
        tf.margin_left = Pt(10)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = lbl
        r.font.name = FONT_BODY
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = WHITE
        if i < len(items) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(10.2), y + Inches(0.55),
                Inches(0.4), Inches(0.22))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEXT_MUTED
            arrow.line.fill.background()
    add_footer(s, 7, total)


def slide_state_management(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "6. Quản lý state - Provider",
                  "Gợi ý bởi nhóm Flutter, dùng ChangeNotifier")

    add_text(s, Inches(0.6), Inches(1.25), Inches(7.5), Inches(0.5),
             "Tại sao chọn Provider?", size=20, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.6), Inches(1.85), Inches(7.5), Inches(3.0), [
        "Đơn giản, hỗ trợ chính thức, học nhanh",
        "Tách UI khỏi logic (separation of concerns)",
        "Có ChangeNotifierProvider, MultiProvider, Consumer, Selector",
        "Kết hợp tốt với async (notifyListeners() sau khi await)",
    ], size=14)

    add_text(s, Inches(0.6), Inches(4.85), Inches(7.5), Inches(0.5),
             "VocabQuest đang dùng 4 Provider", size=20, bold=True,
             color=ACCENT)
    add_bullets(s, Inches(0.6), Inches(5.4), Inches(7.5), Inches(1.8), [
        "UserProvider - thông tin người dùng & đồng bộ Firestore",
        "GameProvider - state của mini-game đang chơi",
        "SettingsProvider - dark mode, ngôn ngữ, âm thanh, nhắc nhẹ",
        "FavoritesProvider - các game / pack yêu thích",
    ], size=14)

    code_box = add_round_rect(s, Inches(8.4), Inches(1.25),
                              Inches(4.4), Inches(5.7), INK, corner=0.04)
    add_text(s, Inches(8.6), Inches(1.4), Inches(4.0), Inches(0.4),
             "main.dart", size=12, bold=True, color=GOLD)
    code = (
        "MultiProvider(\n"
        "  providers: [\n"
        "    ChangeNotifierProvider(\n"
        "      create: (_) => UserProvider()),\n"
        "    ChangeNotifierProvider(\n"
        "      create: (_) =>\n"
        "        SettingsProvider()..loadSettings()),\n"
        "    ChangeNotifierProvider(\n"
        "      create: (_) => GameProvider()),\n"
        "    ChangeNotifierProvider(\n"
        "      create: (_) =>\n"
        "        FavoritesProvider()..load()),\n"
        "  ],\n"
        "  child: Consumer<SettingsProvider>(\n"
        "    builder: (ctx, s, _) =>\n"
        "      MaterialApp(...)\n"
        "  ),\n"
        ");"
    )
    tb = s.shapes.add_textbox(Inches(8.6), Inches(1.85),
                              Inches(4.1), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.name = FONT_MONO
    r.font.size = Pt(11)
    r.font.color.rgb = WHITE
    add_footer(s, 8, total)


def slide_architecture(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "7. Kiến trúc dự án",
                  "Phân lớp rõ ràng theo trách nhiệm")

    layers = [
        ("UI / Screens & Widgets", PRIMARY,
         "screens/, widgets/  -  hiển thị và tương tác"),
        ("State - Providers", ACCENT,
         "user / game / settings / favorites"),
        ("Service - I/O & async",  ACCENT2,
         "firestore_service, auth_service, json_service, audio_service, "
         "notification_service, storage_service, local_storage"),
        ("Model & Config", GOLD,
         "user_model, vocab_model, level_model, game_result_model, "
         "constants, theme, design_tokens"),
        ("External - Firebase, OS APIs", PRIMARY_DARK,
         "Firebase Auth - Firestore - Storage - SharedPreferences - "
         "Notifications - TTS"),
    ]
    base_y = Inches(1.4)
    h = Inches(1.0)
    gap = Inches(0.12)
    for i, (title, col, desc) in enumerate(layers):
        y = base_y + i * (h + gap)
        layer = add_round_rect(s, Inches(0.6), y,
                               Inches(12.1), h, WHITE, corner=0.06)
        bar = add_round_rect(s, Inches(0.6), y,
                             Inches(0.18), h, col, corner=0.5)
        add_text(s, Inches(1.0), y + Inches(0.12),
                 Inches(4.0), Inches(0.4),
                 title, size=15, bold=True, color=col)
        add_text(s, Inches(1.0), y + Inches(0.5),
                 Inches(11.5), Inches(0.5),
                 desc, size=12, color=TEXT_DARK)

    add_footer(s, 9, total)


def slide_libs_firebase(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "8. Thư viện - Backend & Lưu trữ",
                  "Firebase là xương sống dữ liệu")

    cards = [
        ("firebase_core", "^2.24.2", PRIMARY,
         "Khởi tạo & cấu hình Firebase cho mọi nền tảng."),
        ("firebase_auth", "^4.15.3", ACCENT,
         "Đăng ký / đăng nhập email + password, đổi mật khẩu, sign out."),
        ("cloud_firestore", "^4.13.6", ACCENT2,
         "NoSQL realtime: users, game_results, leaderboard. Stream + Transaction."),
        ("firebase_storage", "^11.5.6", GOLD,
         "Upload avatar người dùng từ image_picker."),
        ("shared_preferences", "^2.2.2", PRIMARY_DARK,
         "Lưu setting cục bộ: dark mode, ngôn ngữ, âm thanh, giờ nhắc nhẹ."),
        ("intl / uuid", "0.19 / 4.2", ACCENT,
         "Format ngày giờ, định danh duy nhất cho game_result."),
    ]
    cw = Inches(4.05)
    ch = Inches(1.7)
    gx = Inches(0.6)
    gy = Inches(1.3)
    for i, (name, ver, col, desc) in enumerate(cards):
        col_idx = i % 3
        row_idx = i // 3
        x = gx + col_idx * (cw + Inches(0.1))
        y = gy + row_idx * (ch + Inches(0.25))
        card = add_round_rect(s, x, y, cw, ch, WHITE, corner=0.07)
        bar = add_round_rect(s, x, y, Inches(0.12), ch, col, corner=0.4)
        add_text(s, x + Inches(0.3), y + Inches(0.15),
                 cw - Inches(0.4), Inches(0.4),
                 name, size=15, bold=True, color=col)
        add_text(s, x + Inches(0.3), y + Inches(0.5),
                 cw - Inches(0.4), Inches(0.3),
                 "Phiên bản: " + ver, size=11, color=TEXT_MUTED, italic=True)
        add_text(s, x + Inches(0.3), y + Inches(0.85),
                 cw - Inches(0.5), ch - Inches(1.0),
                 desc, size=12, color=TEXT_DARK)

    add_footer(s, 10, total)


def slide_libs_ui(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "9. Thư viện - UI & Animation",
                  "Tạo trải nghiệm sinh động, hấp dẫn")

    cards = [
        ("provider", "^6.1.1", PRIMARY,
         "Quản lý state theo Provider pattern."),
        ("google_fonts", "^6.1.0", ACCENT,
         "Sử dụng font Google miễn phí ngay trong app."),
        ("font_awesome_flutter", "^10.6", ACCENT2,
         "Bộ icon Font Awesome đầy đủ."),
        ("lucide_icons", "^0.257", PRIMARY_DARK,
         "Bộ icon Lucide hiện đại, mảnh mẽ."),
        ("flutter_animate", "^4.3", GOLD,
         "Tạo animation chuỗi (.fade().slide()...) ngắn gọn."),
        ("lottie", "^2.7.0", PRIMARY,
         "Phát animation Lottie - hiệu ứng ăn mừng, loading."),
        ("confetti", "^0.7.0", ACCENT,
         "Hiệu ứng pháo giấy khi pass level."),
        ("shimmer", "^3.0.0", ACCENT2,
         "Skeleton loading mượt mà."),
        ("animated_text_kit", "^4.2.2", GOLD,
         "Hiệu ứng typing / fade cho text."),
        ("percent_indicator", "^4.2.3", PRIMARY_DARK,
         "Thanh tiến độ tròn, đoạn đường, XP bar."),
        ("cached_network_image", "^3.3", ACCENT,
         "Cache ảnh từ mạng - giảm tải dữ liệu."),
        ("flutter_svg", "^2.0.9", PRIMARY,
         "Render ảnh SVG cho icon vector."),
    ]
    cw = Inches(3.05)
    ch = Inches(1.45)
    gx = Inches(0.45)
    gy = Inches(1.25)
    for i, (name, ver, col, desc) in enumerate(cards):
        col_idx = i % 4
        row_idx = i // 4
        x = gx + col_idx * (cw + Inches(0.1))
        y = gy + row_idx * (ch + Inches(0.2))
        card = add_round_rect(s, x, y, cw, ch, WHITE, corner=0.08)
        bar = add_round_rect(s, x, y, Inches(0.1), ch, col, corner=0.4)
        add_text(s, x + Inches(0.22), y + Inches(0.1),
                 cw - Inches(0.3), Inches(0.35),
                 name, size=12, bold=True, color=col)
        add_text(s, x + Inches(0.22), y + Inches(0.42),
                 cw - Inches(0.3), Inches(0.25),
                 "v" + ver, size=10, color=TEXT_MUTED, italic=True)
        add_text(s, x + Inches(0.22), y + Inches(0.7),
                 cw - Inches(0.3), ch - Inches(0.8),
                 desc, size=10, color=TEXT_DARK)

    add_footer(s, 11, total)


def slide_libs_other(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "10. Thư viện - Audio, Notify, Tiện ích",
                  "Hoàn thiện trải nghiệm người dùng")

    cards = [
        ("audioplayers", "^5.2.1", PRIMARY,
         "Phát hiệu ứng âm thanh: đúng, sai, ăn mừng."),
        ("flutter_tts", "^3.8.5", ACCENT,
         "Text-to-Speech: đọc to từ vựng tiếng Anh giúp luyện phát âm."),
        ("flutter_local_notifications", "^17.2.4", ACCENT2,
         "Nhắc nhở học tập hằng ngày, lập lịch theo giờ người dùng chọn."),
        ("timezone", "^0.9.4", GOLD,
         "Đặt lịch notification chính xác theo múi giờ thiết bị."),
        ("image_picker", "^1.0.7", PRIMARY_DARK,
         "Chọn ảnh từ camera/galery để làm avatar."),
        ("flutter_launcher_icons", "^0.14", ACCENT,
         "Sinh icon ứng dụng từ logo.png cho mọi nền tảng."),
    ]
    cw = Inches(4.05)
    ch = Inches(1.7)
    gx = Inches(0.6)
    gy = Inches(1.3)
    for i, (name, ver, col, desc) in enumerate(cards):
        col_idx = i % 3
        row_idx = i // 3
        x = gx + col_idx * (cw + Inches(0.1))
        y = gy + row_idx * (ch + Inches(0.25))
        card = add_round_rect(s, x, y, cw, ch, WHITE, corner=0.07)
        bar = add_round_rect(s, x, y, Inches(0.12), ch, col, corner=0.4)
        add_text(s, x + Inches(0.3), y + Inches(0.15),
                 cw - Inches(0.4), Inches(0.4),
                 name, size=15, bold=True, color=col)
        add_text(s, x + Inches(0.3), y + Inches(0.5),
                 cw - Inches(0.4), Inches(0.3),
                 "Phiên bản: " + ver, size=11, color=TEXT_MUTED, italic=True)
        add_text(s, x + Inches(0.3), y + Inches(0.85),
                 cw - Inches(0.5), ch - Inches(1.0),
                 desc, size=12, color=TEXT_DARK)

    add_footer(s, 12, total)


def slide_use_case_a(prs, total):
    """Phân tích chức năng - Phần 1: Tài khoản & Học chơi."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "11. Phân tích chức năng (1/2) - Tài khoản & Học chơi",
                  "Use-case chính: Auth, Hồ sơ, Vào trận, Phát âm")

    groups = [
        ("Quản lý tài khoản", PRIMARY, [
            "Đăng ký bằng email + mật khẩu (Firebase Auth)",
            "Validate mật khẩu >= 6 ký tự, email đúng định dạng",
            "Đăng nhập, ghi nhớ phiên với Firebase persistence",
            "Quên mật khẩu - gửi email reset",
            "Đăng xuất - xóa session local + Firestore stream",
            "Tạo doc users/{uid} + thưởng 100 coin khởi điểm",
        ]),
        ("Hồ sơ cá nhân", ACCENT, [
            "Xem thông tin: tên, email, level, XP, coin, streak, tim",
            "Sửa tên hiển thị (cập nhật realtime Firestore)",
            "Đổi avatar: image_picker -> resize 400x400 -> base64 lưu Firestore",
            "Đổi mật khẩu (yêu cầu mật khẩu cũ)",
            "Xem lịch sử các lần chơi gần nhất",
        ]),
        ("Chọn gói & vào trận", ACCENT2, [
            "3 gói từ vựng: Sơ cấp / Trung cấp / Nâng cao",
            "Mỗi gói chia 3 level (mức khó tăng dần)",
            "Bản đồ level dạng stage: chưa mở, đang mở, đã pass (1-3 sao)",
            "Phải đạt >= 2 sao mới mở được level kế tiếp",
            "Chọn 1 trong 4 mini-game cho mỗi level",
        ]),
        ("Phát âm từ vựng (TTS)", GOLD, [
            "Tích hợp flutter_tts với giọng en-US, speech rate 0.5",
            "Trong Quiz: bấm icon loa nghe từ tiếng Anh",
            "Trong Matching: chạm vào ô từ để nghe phát âm",
            "Hỗ trợ luyện kỹ năng nghe + phát âm song song",
        ]),
    ]
    cw = Inches(6.0)
    ch = Inches(2.85)
    gx = Inches(0.55)
    gy = Inches(1.3)
    for i, (title, col, items) in enumerate(groups):
        col_idx = i % 2
        row_idx = i // 2
        x = gx + col_idx * (cw + Inches(0.2))
        y = gy + row_idx * (ch + Inches(0.2))
        card = add_round_rect(s, x, y, cw, ch, WHITE, corner=0.06)
        bar = add_round_rect(s, x, y, Inches(0.15), ch, col, corner=0.5)
        add_text(s, x + Inches(0.35), y + Inches(0.15),
                 cw - Inches(0.5), Inches(0.4),
                 title, size=16, bold=True, color=col)
        add_bullets(s, x + Inches(0.35), y + Inches(0.7),
                    cw - Inches(0.5), ch - Inches(0.8),
                    items, size=11, bullet_color=col, line_spacing=1.15)

    add_footer(s, 13, total)


def slide_use_case_b(prs, total):
    """Phân tích chức năng - Phần 2: Tiến độ, Thưởng, Cá nhân hóa, Khác."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "12. Phân tích chức năng (2/2) - Phần thưởng & Cá nhân hóa",
                  "Use-case: Gamification, Leaderboard, Settings, Shop")

    groups = [
        ("Hệ thống XP & Level", PRIMARY, [
            "10 mức level từ 0 -> 12000 XP (theo LevelSystem.xpThresholds)",
            "Mỗi câu đúng cộng score -> quy đổi 1:1 sang XP cơ bản",
            "Bonus accuracy 100% : +30 XP, 80%+ : +15 XP",
            "Bonus level: Trung cấp x1.5, Nâng cao x2.0 XP",
            "Lên level hiển thị animation chúc mừng",
        ]),
        ("Coin & Tim & Cửa hàng", ACCENT, [
            "Coin = round(score x 0.5) + bonus accuracy + bonus speed (<30s)",
            "Bắt đầu: 100 coin, 5 tim. Sai trong Quiz mất 1 tim",
            "Coin dùng để: gợi ý Word Puzzle (10 coin/lần)",
            "Pass level lần đầu nhận thưởng: L1=30c/20XP, L2=60c/40XP, L3=150c/100XP",
            "Shop: mua thêm gói từ vựng nâng cao",
        ]),
        ("Streak & Bảng xếp hạng", ACCENT2, [
            "Streak: số ngày liên tiếp chơi ít nhất 1 trận",
            "Đứt streak nếu nghỉ quá 1 ngày, lưu longestStreak kỷ lục",
            "Bonus streak coin/XP khi đạt mốc 3, 7, 14, 30 ngày",
            "Leaderboard realtime qua Stream Firestore, xếp theo totalXP",
            "Hiển thị top 50 + thứ hạng của user hiện tại",
        ]),
        ("Cá nhân hóa & Tiện ích", GOLD, [
            "Dark mode / Light mode (lưu shared_preferences)",
            "Song ngữ Việt - Anh, đổi ngay không cần restart",
            "Bật/tắt SFX, nhạc nền, hiệu ứng hoạt hình",
            "Nhắc nhở học tập: chọn giờ - flutter_local_notifications",
            "Yêu thích game / pack - hiển thị nhanh ở Home",
            "Lịch sử 50 trận gần nhất, lọc theo loại game",
        ]),
    ]
    cw = Inches(6.0)
    ch = Inches(2.85)
    gx = Inches(0.55)
    gy = Inches(1.3)
    for i, (title, col, items) in enumerate(groups):
        col_idx = i % 2
        row_idx = i // 2
        x = gx + col_idx * (cw + Inches(0.2))
        y = gy + row_idx * (ch + Inches(0.2))
        card = add_round_rect(s, x, y, cw, ch, WHITE, corner=0.06)
        bar = add_round_rect(s, x, y, Inches(0.15), ch, col, corner=0.5)
        add_text(s, x + Inches(0.35), y + Inches(0.15),
                 cw - Inches(0.5), Inches(0.4),
                 title, size=16, bold=True, color=col)
        add_bullets(s, x + Inches(0.35), y + Inches(0.7),
                    cw - Inches(0.5), ch - Inches(0.8),
                    items, size=11, bullet_color=col, line_spacing=1.15)

    add_footer(s, 14, total)


def slide_scoring_overview(prs, total):
    """Tổng quan luật tính điểm chung cho 4 game."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "13. Luật tính điểm & Phần thưởng chung",
                  "Áp dụng cho cả 4 mini-game")

    add_text(s, Inches(0.55), Inches(1.2), Inches(7.5), Inches(0.5),
             "Quy đổi điểm -> XP & Coin", size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.55), Inches(1.8), Inches(7.5), Inches(2.6), [
        "XP nhận = score (1 điểm = 1 XP cơ bản)",
        "Coin nhận = round(score x 0.5)",
        "Tính accuracy = correctAnswers / totalQuestions",
        "Bonus accuracy 100% : +50 coin, +30 XP",
        "Bonus accuracy >= 80% : +20 coin, +15 XP",
        "Bonus tốc độ (< 30 giây): +10 coin",
    ], size=12, line_spacing=1.2)

    add_text(s, Inches(0.55), Inches(4.6), Inches(7.5), Inches(0.5),
             "Hệ số nhân theo độ khó", size=18, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.55), Inches(5.2), Inches(7.5), Inches(2.0), [
        "Sơ cấp (beginner): hệ số XP x 1.0",
        "Trung cấp (intermediate): hệ số XP x 1.5",
        "Nâng cao (advanced): hệ số XP x 2.0",
        "Coin không nhân theo level",
    ], size=12, line_spacing=1.2)

    box = add_round_rect(s, Inches(8.4), Inches(1.2),
                         Inches(4.4), Inches(5.7), WHITE, corner=0.06)
    add_text(s, Inches(8.65), Inches(1.35), Inches(4.0), Inches(0.4),
             "Quy ước Sao (Stars)", size=15, bold=True, color=PRIMARY)
    add_text(s, Inches(8.65), Inches(1.8), Inches(4.0), Inches(0.4),
             "Dựa trên độ chính xác:", size=11, color=TEXT_MUTED, italic=True)

    star_rules = [
        ("3 sao", ">= 90% đúng", GOLD),
        ("2 sao", ">= 60% đúng", ACCENT),
        ("1 sao", ">= 30% đúng", PRIMARY),
        ("0 sao", "< 30% (chưa pass)", TEXT_MUTED),
    ]
    for i, (lbl, cond, col) in enumerate(star_rules):
        y = Inches(2.3) + i * Inches(0.7)
        chip = add_round_rect(s, Inches(8.65), y,
                              Inches(3.9), Inches(0.55), col, corner=0.4)
        tf = chip.text_frame
        tf.margin_left = Pt(10)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = lbl + "  -  "
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        r1.font.size = Pt(13)
        r1.font.name = FONT_TITLE
        r2 = p.add_run()
        r2.text = cond
        r2.font.color.rgb = WHITE
        r2.font.size = Pt(12)
        r2.font.name = FONT_BODY

    note = add_round_rect(s, Inches(8.65), Inches(5.3),
                          Inches(3.9), Inches(1.4), BG_LIGHT, corner=0.1)
    add_text(s, Inches(8.85), Inches(5.45), Inches(3.6), Inches(0.4),
             "Lưu ý mở khóa", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(8.85), Inches(5.75), Inches(3.6), Inches(1.0),
             "Phải đạt tối thiểu 2 sao thì level kế tiếp mới được mở "
             "khóa và thưởng coin/XP mới được ghi nhận lần đầu.",
             size=11, color=TEXT_DARK)

    add_footer(s, 15, total)


def slide_data_flow(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "19. Luồng dữ liệu - Chơi 1 game",
                  "Kịch bản điển hình từ UI -> Provider -> Firestore")

    steps = [
        ("UI: Người dùng chọn gói từ, bấm Start",  PRIMARY),
        ("GameProvider.loadVocab(level)  -> JsonService đọc JSON cục bộ",
         ACCENT),
        ("Chơi game: cập nhật score / correct / wrong qua notifyListeners()",
         ACCENT2),
        ("Kết thúc: GameProvider.finishGame() tính điểm, XP, coin, streak",
         GOLD),
        ("FirestoreService.runTransaction() ghi atomic vào users + game_results",
         PRIMARY_DARK),
        ("UserProvider lắng nghe Stream<User> -> UI tự cập nhật XP / level",
         ACCENT),
    ]
    for i, (txt, col) in enumerate(steps):
        y = Inches(1.4) + i * Inches(0.85)
        c = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.65), Inches(0.65))
        c.fill.solid()
        c.fill.fore_color.rgb = col
        c.line.fill.background()
        tf = c.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.bold = True
        r.font.size = Pt(20)
        r.font.color.rgb = WHITE
        r.font.name = FONT_TITLE

        chip = add_round_rect(s, Inches(1.5), y + Inches(0.05),
                              Inches(11.2), Inches(0.55), WHITE, corner=0.4)
        add_text(s, Inches(1.7), y + Inches(0.08),
                 Inches(10.8), Inches(0.5),
                 txt, size=14, bold=True, color=TEXT_DARK)

    add_footer(s, 21, total)


def slide_db_schema(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "20. Sơ đồ dữ liệu - Firestore",
                  "3 collection chính + cấu hình local")

    add_text(s, Inches(0.6), Inches(1.2), Inches(7.0), Inches(0.5),
             "Firestore Collections", size=18, bold=True, color=PRIMARY)

    cols = [
        ("users / {uid}", PRIMARY,
         ["uid, email, displayName, avatarUrl",
          "totalScore, totalCoins, totalXP, level",
          "streak, longestStreak, hearts",
          "ownedPacks: List<String>",
          "progress: Map<gameType|packId, level>",
          "lastPlayedDate, createdAt"]),
        ("game_results / {auto}", ACCENT,
         ["uid, gameType, level (pack)",
          "score, correctAnswers, totalQuestions",
          "duration (seconds)",
          "coinsEarned, xpEarned",
          "playedAt: Timestamp"]),
        ("user_settings (local)", ACCENT2,
         ["isDarkMode, language",
          "soundEnabled, musicEnabled",
          "notificationEnabled",
          "reminderHour, reminderMinute",
          "favoriteGames: List<String>",
          "Lưu bằng shared_preferences"]),
    ]
    cw = Inches(4.05)
    ch = Inches(4.8)
    gx = Inches(0.55)
    gy = Inches(1.8)
    for i, (title, col, fields) in enumerate(cols):
        x = gx + i * (cw + Inches(0.1))
        card = add_round_rect(s, x, gy, cw, ch, WHITE, corner=0.06)
        head = add_round_rect(s, x, gy, cw, Inches(0.7), col, corner=0.06)
        add_text(s, x + Inches(0.3), gy + Inches(0.15),
                 cw - Inches(0.4), Inches(0.4),
                 title, size=14, bold=True, color=WHITE,
                 font=FONT_MONO)
        for j, f in enumerate(fields):
            y = gy + Inches(0.85) + j * Inches(0.55)
            line = add_round_rect(s, x + Inches(0.2), y,
                                  cw - Inches(0.4), Inches(0.45),
                                  BG_LIGHT, corner=0.4)
            tf = line.text_frame
            tf.margin_left = Pt(8)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = f
            r.font.name = FONT_MONO
            r.font.size = Pt(11)
            r.font.color.rgb = TEXT_DARK

    add_footer(s, 22, total)


def slide_screens_overview(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "21. Các màn hình chính",
                  "Tổng quan luồng UI của VocabQuest")

    items = [
        ("Splash", "Khởi tạo Firebase, audio, kiểm tra đăng nhập", PRIMARY),
        ("Auth", "Login / Register với validation", ACCENT),
        ("Home", "Hero greeting, stats, action tile", ACCENT2),
        ("Game Menu", "Featured + grid 4 mini-game", GOLD),
        ("Pack Selection", "Chọn gói từ vựng theo level", PRIMARY_DARK),
        ("Level Map", "Bản đồ level dạng stage, hiện sao", PRIMARY),
        ("Mini-game", "Matching / Quiz / Puzzle / Memory", ACCENT),
        ("Game Result", "Điểm - sao - thưởng + confetti", ACCENT2),
        ("Profile", "Thông tin, edit, đổi mật khẩu", GOLD),
        ("Leaderboard", "BXH realtime theo XP", PRIMARY_DARK),
        ("History", "Lịch sử chơi game", PRIMARY),
        ("Shop", "Mua gói từ vựng bằng coin", ACCENT),
        ("Favorites", "Game / pack yêu thích", ACCENT2),
        ("Settings", "Theme, ngôn ngữ, âm thanh, nhắc nhở", GOLD),
    ]
    cw = Inches(2.95)
    ch = Inches(1.0)
    gx = Inches(0.55)
    gy = Inches(1.25)
    for i, (name, desc, col) in enumerate(items):
        col_idx = i % 4
        row_idx = i // 4
        x = gx + col_idx * (cw + Inches(0.1))
        y = gy + row_idx * (ch + Inches(0.15))
        card = add_round_rect(s, x, y, cw, ch, WHITE, corner=0.08)
        head = add_round_rect(s, x, y, Inches(0.1), ch, col, corner=0.5)
        add_text(s, x + Inches(0.2), y + Inches(0.08),
                 cw - Inches(0.3), Inches(0.35),
                 name, size=13, bold=True, color=col)
        add_text(s, x + Inches(0.2), y + Inches(0.42),
                 cw - Inches(0.3), ch - Inches(0.5),
                 desc, size=10, color=TEXT_DARK)

    add_footer(s, 23, total)


def _draw_game_detail(slide, *, idx, title, badge, color, intro,
                      rules, scoring, end_cond, footer_no, total):
    """Helper vẽ slide chi tiết 1 mini-game."""
    set_solid_bg(slide, BG_LIGHT)
    add_title_bar(slide, f"{idx}. Mini-game - {title}",
                  "Cách chơi - Luật tính điểm - Điều kiện kết thúc")

    # Header chip name + badge
    head = add_round_rect(slide, Inches(0.55), Inches(1.2),
                          Inches(12.2), Inches(0.85), color, corner=0.3)
    add_text(slide, Inches(0.85), Inches(1.32),
             Inches(7.0), Inches(0.6),
             title, font=FONT_TITLE, size=24, bold=True, color=WHITE)
    b = add_round_rect(slide, Inches(11.0), Inches(1.45),
                       Inches(1.6), Inches(0.4), WHITE, corner=0.5)
    tf = b.text_frame
    tf.margin_left = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = badge
    r.font.bold = True
    r.font.size = Pt(11)
    r.font.color.rgb = color
    r.font.name = FONT_TITLE
    add_text(slide, Inches(0.85), Inches(1.7),
             Inches(11.0), Inches(0.4),
             intro, font=FONT_BODY, size=12, color=WHITE, italic=True)

    # 3 cột: Luật chơi - Tính điểm - Kết thúc
    col_w = Inches(4.0)
    col_h = Inches(4.85)
    gap = Inches(0.1)
    x0 = Inches(0.55)
    y0 = Inches(2.25)

    sections = [
        ("Luật chơi", PRIMARY, rules),
        ("Tính điểm", ACCENT, scoring),
        ("Kết thúc & Thưởng", ACCENT2, end_cond),
    ]
    for i, (sec_title, sec_col, items) in enumerate(sections):
        x = x0 + i * (col_w + gap)
        card = add_round_rect(slide, x, y0, col_w, col_h, WHITE, corner=0.06)
        head_card = add_round_rect(slide, x, y0,
                                   col_w, Inches(0.55), sec_col, corner=0.3)
        add_text(slide, x + Inches(0.25), y0 + Inches(0.12),
                 col_w - Inches(0.4), Inches(0.4),
                 sec_title, font=FONT_TITLE, size=15, bold=True, color=WHITE)
        add_bullets(slide, x + Inches(0.2), y0 + Inches(0.7),
                    col_w - Inches(0.3), col_h - Inches(0.8),
                    items, size=11, bullet_color=sec_col, line_spacing=1.2)

    add_footer(slide, footer_no, total)


def slide_minigames_overview(prs, total):
    """Slide 16 - Tổng quan 4 mini-game."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "14. Tổng quan 4 mini-game",
                  "Mỗi game rèn 1 kỹ năng từ vựng khác nhau")

    games = [
        ("Nối từ (Matching)", PRIMARY, "Hot",
         "Kéo - chạm nối từ tiếng Anh với nghĩa tiếng Việt",
         "Kỹ năng: nhận diện - ghi nhớ ngữ nghĩa",
         "4 cặp/round  -  3 round  -  60s"),
        ("Trắc nghiệm (Quiz)", ACCENT, "Cơ bản",
         "Hiển thị từ tiếng Anh, chọn nghĩa đúng trong 4 đáp án",
         "Kỹ năng: phản xạ + nghe TTS",
         "10 câu  -  15s/câu"),
        ("Xếp chữ (Word Puzzle)", ACCENT2, "Thử thách",
         "Sắp xếp các chữ cái xáo trộn thành từ tiếng Anh đúng",
         "Kỹ năng: chính tả - đánh vần",
         "5 từ/lượt  -  có gợi ý"),
        ("Lật thẻ (Memory)", GOLD, "Mới",
         "Lật 2 thẻ tìm cặp từ - nghĩa giống nhau (Pexeso)",
         "Kỹ năng: trí nhớ ngắn hạn + cặp đôi",
         "6 cặp/round  -  3 round  -  90s"),
    ]
    cw = Inches(6.0)
    ch = Inches(2.85)
    gx = Inches(0.55)
    gy = Inches(1.3)
    for i, (name, col, badge, desc, skill, time_info) in enumerate(games):
        col_idx = i % 2
        row_idx = i // 2
        x = gx + col_idx * (cw + Inches(0.2))
        y = gy + row_idx * (ch + Inches(0.2))
        card = add_round_rect(s, x, y, cw, ch, WHITE, corner=0.07)
        bar = add_round_rect(s, x, y, Inches(0.18), ch, col, corner=0.5)

        add_text(s, x + Inches(0.4), y + Inches(0.18),
                 Inches(4.0), Inches(0.5),
                 name, size=20, bold=True, color=col)

        b = add_round_rect(s, x + cw - Inches(1.5), y + Inches(0.25),
                           Inches(1.2), Inches(0.4), col, corner=0.5)
        tf = b.text_frame
        tf.margin_left = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = badge
        r.font.bold = True
        r.font.size = Pt(11)
        r.font.color.rgb = WHITE
        r.font.name = FONT_TITLE

        add_text(s, x + Inches(0.4), y + Inches(0.8),
                 cw - Inches(0.6), Inches(0.45),
                 desc, size=12, color=TEXT_DARK)
        add_text(s, x + Inches(0.4), y + Inches(1.3),
                 cw - Inches(0.6), Inches(0.4),
                 skill, size=11, color=TEXT_MUTED, italic=True)

        time_chip = add_round_rect(s, x + Inches(0.4), y + Inches(1.85),
                                   Inches(4.0), Inches(0.45),
                                   BG_LIGHT, corner=0.4)
        add_text(s, x + Inches(0.6), y + Inches(1.92),
                 Inches(4.0), Inches(0.35),
                 time_info, size=11, bold=True, color=col)

    add_footer(s, 16, total)


def slide_game_matching(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_game_detail(
        s,
        idx="15",
        title="Nối từ (Matching)",
        badge="Hot",
        color=PRIMARY,
        intro="Kéo - chạm nối các cặp từ tiếng Anh và nghĩa tiếng Việt "
              "trong thời gian giới hạn",
        rules=[
            "Mỗi round hiển thị 4 cặp (8 ô): cột trái là từ EN, cột phải nghĩa VN",
            "Chạm 1 ô EN, sau đó chạm 1 ô VN để thử nối",
            "Đúng: 2 ô đổi xanh, biến mất khỏi lưới (matched)",
            "Sai: 2 ô nháy đỏ trong 600ms rồi tự bỏ chọn",
            "Có thể bỏ chọn ô đã chọn trước khi chốt cặp",
            "Round mới khi đã match hết 4 cặp, tổng 3 round/trận",
            "Pool dedupe: không lặp từ giữa các round trong cùng trận",
        ],
        scoring=[
            "+10 điểm cho mỗi cặp nối đúng LẦN ĐẦU",
            "Cặp đã từng sai (failed) sẽ KHÔNG được tính điểm dù sau nối đúng",
            "Sai 1 lần đánh dấu pair là failed: +1 wrong vào totalQuestions",
            "Nối đúng pair chưa failed: +1 correct và +1 totalQuestions",
            "Tối đa 4 pairs x 3 rounds = 12 cặp -> 120 điểm",
            "Audio: playCorrect khi đúng, playWrong khi sai",
        ],
        end_cond=[
            "Hết giờ 60 giây hoặc match hết 12 cặp -> kết thúc",
            "Tính accuracy = correct / total -> 1-3 sao",
            "Quy ra XP, coin theo công thức chung (slide 12)",
            "Pass >=2 sao + chưa pass trước đó -> mở level kế tiếp + thưởng",
            "Cập nhật streak, longestStreak, ghi vào game_results",
        ],
        footer_no=17, total=total,
    )


def slide_game_quiz(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_game_detail(
        s,
        idx="16",
        title="Trắc nghiệm (Quiz)",
        badge="Cơ bản",
        color=ACCENT,
        intro="Hiển thị 1 từ tiếng Anh, chọn nghĩa đúng trong 4 đáp án "
              "với áp lực thời gian từng câu",
        rules=[
            "Mỗi trận lấy ngẫu nhiên 10 từ từ pool theo level",
            "Mỗi câu sinh 4 đáp án: 1 đúng + 3 distractor random từ pool",
            "Có icon loa (TTS) phát âm tiếng Anh từ hỏi",
            "Người chơi chạm 1 đáp án để chốt -> hiện màu đúng/sai",
            "Hết 1.5s sẽ tự sang câu kế",
            "Mỗi câu giới hạn 15 giây (đếm ngược thanh ngang)",
            "Không trả lời kịp = sai (timeout)",
        ],
        scoring=[
            "Trả lời đúng: +(10 + thời gian còn lại) điểm",
            "Ví dụ: đúng khi còn 12s -> +22 điểm",
            "Khuyến khích phản xạ nhanh: tối đa 25/câu",
            "Trả lời sai hoặc timeout: 0 điểm + 1 wrong",
            "10 câu x 25 = tối đa 250 điểm",
            "Audio đúng/sai phát ngay khi chốt đáp án",
        ],
        end_cond=[
            "Kết thúc khi đã qua hết 10 câu",
            "Accuracy = correct / 10 -> 1-3 sao",
            "Quy đổi XP và coin theo công thức chung",
            "Bonus accuracy 100% (+30 XP, +50 coin) khá khả thi với Quiz",
            "Pass >=2 sao mới mở level kế tiếp + thưởng level lần đầu",
        ],
        footer_no=18, total=total,
    )


def slide_game_puzzle(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_game_detail(
        s,
        idx="17",
        title="Xếp chữ (Word Puzzle)",
        badge="Thử thách",
        color=ACCENT2,
        intro="Sắp xếp các chữ cái bị xáo trộn để tạo thành từ tiếng Anh "
              "đúng, có hỗ trợ gợi ý",
        rules=[
            "Mỗi trận chọn 5 từ ngẫu nhiên theo level",
            "Hệ thống xáo trộn các chữ cái và cho hiện ở khu chứa dưới",
            "Người chơi chạm chữ cái để đẩy lên ô trả lời, theo thứ tự",
            "Có thể chạm lại ô đã điền để bỏ ra (sửa thứ tự)",
            "Khi điền đủ độ dài từ -> tự kiểm tra đúng/sai",
            "Có nút gợi ý (hint): điền chữ kế tiếp đúng tự động",
            "Mỗi lần dùng hint trừ 10 coin (cần >=10 coin)",
        ],
        scoring=[
            "Trả lời đúng: +20 điểm",
            "Trả lời sai: 0 điểm + 1 wrong",
            "Mỗi câu có nhiều cơ hội điền - chỉ chốt khi đủ chữ cái",
            "5 từ x 20 = tối đa 100 điểm",
            "Tip: dùng hint thông minh - mất coin nhưng đảm bảo đúng",
            "Sau 1.8s tự chuyển từ kế tiếp",
        ],
        end_cond=[
            "Kết thúc khi qua hết 5 từ",
            "Accuracy = correct / 5 -> 1-3 sao",
            "Quy đổi XP, coin và bonus theo công thức chung",
            "Pass >=2 sao mở level kế + thưởng coin/XP lần đầu",
            "Lưu game_result, cập nhật streak realtime",
        ],
        footer_no=19, total=total,
    )


def slide_game_memory(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_game_detail(
        s,
        idx="18",
        title="Lật thẻ (Memory / Pexeso)",
        badge="Mới",
        color=GOLD,
        intro="Lật 2 thẻ tìm cặp từ tiếng Anh - nghĩa tiếng Việt giống "
              "nhau, rèn trí nhớ ngắn hạn",
        rules=[
            "Mỗi round hiển thị grid 4x3 = 12 thẻ úp (6 cặp)",
            "Chạm 1 thẻ -> lật mở. Chạm thẻ thứ 2 -> kiểm tra match",
            "Thẻ chứa từ EN hoặc nghĩa VN của cùng vocabId là 1 cặp",
            "Khớp: cả 2 thẻ giữ mở (matched), cấm chạm tiếp",
            "Không khớp: úp lại sau 800ms, có thể nhớ vị trí",
            "Round mới khi đủ 6 cặp matched, tổng 3 round/trận",
            "Pool dedupe: không lặp cặp giữa các round",
        ],
        scoring=[
            "+20 điểm cho mỗi cặp khớp LẦN ĐẦU",
            "Cặp đã từng mismatch (failed) sẽ KHÔNG cộng điểm dù sau khớp",
            "Mismatch 1 lần đánh dấu pair là failed: +1 wrong",
            "Khớp pair chưa failed: +1 correct và +1 totalQuestions",
            "Tối đa 6 pairs x 3 rounds = 18 cặp -> 360 điểm",
            "Game có điểm cao nhất trong 4 mini-game",
        ],
        end_cond=[
            "Hết giờ 90 giây hoặc match hết 18 cặp -> kết thúc",
            "Accuracy = correct / total -> 1-3 sao",
            "Quy đổi XP, coin theo công thức chung (slide 12)",
            "Bonus speed nếu hoàn thành dưới 30s (+10 coin)",
            "Pass >=2 sao mở level kế tiếp + thưởng lần đầu",
        ],
        footer_no=20, total=total,
    )


def slide_async_in_app(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "22. Bất đồng bộ - Ứng dụng trong VocabQuest",
                  "Future, async/await, Stream cụ thể trong code")

    add_text(s, Inches(0.55), Inches(1.2), Inches(7.5), Inches(0.5),
             "Điểm chốt trong code", size=20, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.55), Inches(1.85), Inches(7.5), Inches(5.0), [
        "main(): Future<void> async - khởi tạo Firebase, audio, "
        "notification trước khi runApp",
        "JsonService.loadVocab(): Future - đọc file JSON từ rootBundle",
        "FirestoreService.getUser(): Future - đọc 1 lần",
        "FirestoreService.streamUser(): Stream - lắng nghe realtime",
        "FirestoreService.runTransaction(): bảo toàn atomic khi cập nhật "
        "progress + XP + coin",
        "GameProvider.finishGame(): gọi nhiều await tuần tự, sau đó "
        "notifyListeners() cập nhật UI",
        "AudioService / NotificationService init() gọi trong main bằng "
        "try - catch để lỗi không chặn app khởi động",
    ], size=13)

    code_box = add_round_rect(s, Inches(8.4), Inches(1.2),
                              Inches(4.4), Inches(5.7), INK, corner=0.04)
    add_text(s, Inches(8.6), Inches(1.35), Inches(4.0), Inches(0.4),
             "firestore_service.dart", size=11, bold=True, color=GOLD)
    code = (
        "Stream<UserModel?> streamUser(String uid){\n"
        "  return _db.collection('users')\n"
        "    .doc(uid).snapshots()\n"
        "    .map((doc) =>\n"
        "      doc.exists\n"
        "        ? UserModel.fromFirestore(doc)\n"
        "        : null);\n"
        "}\n\n"
        "Future<bool> updateLevelProgress(...) async {\n"
        "  return _db.runTransaction<bool>((txn) async {\n"
        "    final snap = await txn.get(docRef);\n"
        "    if (level <= current) return false;\n"
        "    txn.update(docRef, {\n"
        "      'progress.\\$key': level,\n"
        "      'totalCoins':\n"
        "         FieldValue.increment(coin),\n"
        "      'totalXP':\n"
        "         FieldValue.increment(xp),\n"
        "    });\n"
        "    return true;\n"
        "  });\n"
        "}"
    )
    tb = s.shapes.add_textbox(Inches(8.6), Inches(1.85),
                              Inches(4.1), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.name = FONT_MONO
    r.font.size = Pt(10)
    r.font.color.rgb = WHITE
    add_footer(s, 24, total)


def slide_done(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, (0xE8, 0xFA, 0xF0), (0xFF, 0xFF, 0xFF), angle=160)
    add_title_bar(s, "23. Tính năng ĐÃ HOÀN THÀNH",
                  "Phần bản đầu tư siêu - sẵn sàng demo")

    rows = [
        ("Xác thực người dùng",
         "Firebase Auth: đăng ký, đăng nhập, đổi mật khẩu, đăng xuất",
         ACCENT2),
        ("4 mini-game đầy đủ",
         "Matching, Quiz, Word Puzzle, Memory - chạy tốt trên Android",
         PRIMARY),
        ("Hệ thống tiến độ",
         "Điểm - XP - Level - Coin - Tim - Streak - Longest streak",
         ACCENT),
        ("Bảng xếp hạng realtime",
         "Stream Firestore, sắp xếp theo totalXP, hiện top + thứ hạng user",
         GOLD),
        ("Lịch sử và thống kê",
         "Game results theo thời gian, biểu đồ tiến độ",
         PRIMARY_DARK),
        ("Hồ sơ & avatar",
         "Edit Profile, upload avatar bằng image_picker + Firebase Storage",
         ACCENT2),
        ("Cá nhân hóa",
         "Dark/Light, song ngữ Việt/Anh, âm thanh, nhạc nền, nhắc nhở",
         PRIMARY),
        ("Phát âm từ vựng",
         "Tích hợp flutter_tts - đọc to từ khi chọn ở Quiz/Matching",
         ACCENT),
    ]
    cw = Inches(6.05)
    ch = Inches(1.4)
    gx = Inches(0.55)
    gy = Inches(1.2)
    for i, (title, desc, col) in enumerate(rows):
        col_idx = i % 2
        row_idx = i // 2
        x = gx + col_idx * (cw + Inches(0.15))
        y = gy + row_idx * (ch + Inches(0.15))
        card = add_round_rect(s, x, y, cw, ch, WHITE, corner=0.06)
        chk = s.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.3),
            Inches(0.7), Inches(0.7))
        chk.fill.solid()
        chk.fill.fore_color.rgb = col
        chk.line.fill.background()
        tf = chk.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "OK"
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.size = Pt(13)
        r.font.name = FONT_TITLE

        add_text(s, x + Inches(1.0), y + Inches(0.2),
                 cw - Inches(1.2), Inches(0.4),
                 title, size=14, bold=True, color=col)
        add_text(s, x + Inches(1.0), y + Inches(0.6),
                 cw - Inches(1.2), ch - Inches(0.7),
                 desc, size=11, color=TEXT_DARK)

    add_footer(s, 25, total)


def slide_undone(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, (0xFF, 0xF3, 0xE0), (0xFF, 0xFF, 0xFF), angle=160)
    add_title_bar(s, "24. Hạn chế & chưa hoàn thiện",
                  "Trung thực - hiểu để cải thiện")

    rows = [
        ("Chưa kiểm thử tự động",
         "Chưa viết unit/widget test cho các game và Firestore service.",
         ACCENT),
        ("Chưa đăng nhập mạng xã hội",
         "Chưa tích hợp Google / Apple Sign-In - chỉ có email-password.",
         PRIMARY),
        ("Gói từ vựng còn hạn chế",
         "Dữ liệu nằm trong file JSON, chưa có backend cập nhật động.",
         ACCENT2),
        ("Chưa có chế độ offline đầy đủ",
         "Mất mạng có thể không nạp được avatar, BXH; chưa có cache 100%.",
         PRIMARY_DARK),
        ("Chưa có tính năng xã hội",
         "Chưa có kết bạn, thách đấu, chia sẻ kết quả.",
         GOLD),
        ("Chưa phát hành iOS",
         "Đã cấu hình nhưng chưa test kỹ trên iOS, chưa publish App Store.",
         ACCENT),
        ("Chưa tích hợp AI",
         "Chưa có chatbot luyện tập, chưa có adaptive learning theo skill.",
         PRIMARY),
        ("Hiệu năng với data lớn",
         "BXH chỉ load top 50, chưa phân trang khi user vượt 1000.",
         ACCENT2),
    ]
    cw = Inches(6.05)
    ch = Inches(1.4)
    gx = Inches(0.55)
    gy = Inches(1.2)
    for i, (title, desc, col) in enumerate(rows):
        col_idx = i % 2
        row_idx = i // 2
        x = gx + col_idx * (cw + Inches(0.15))
        y = gy + row_idx * (ch + Inches(0.15))
        card = add_round_rect(s, x, y, cw, ch, WHITE, corner=0.06)
        chk = s.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.3),
            Inches(0.7), Inches(0.7))
        chk.fill.solid()
        chk.fill.fore_color.rgb = col
        chk.line.fill.background()
        tf = chk.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "!"
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.size = Pt(18)
        r.font.name = FONT_TITLE

        add_text(s, x + Inches(1.0), y + Inches(0.2),
                 cw - Inches(1.2), Inches(0.4),
                 title, size=14, bold=True, color=col)
        add_text(s, x + Inches(1.0), y + Inches(0.6),
                 cw - Inches(1.2), ch - Inches(0.7),
                 desc, size=11, color=TEXT_DARK)

    add_footer(s, 26, total)


def slide_lessons(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "25. Khó khăn & bài học",
                  "Từ triển khai thực tế")

    add_text(s, Inches(0.6), Inches(1.25), Inches(5.8), Inches(0.5),
             "Khó khăn gặp phải", size=20, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.6), Inches(1.85), Inches(5.8), Inches(5.0), [
        "Ghi atomic vào Firestore: phải dùng Transaction để tránh "
        "race-condition khi cộng XP/coin",
        "Quản lý state khi có async: phải gọi notifyListeners() đúng lúc, "
        "tránh setState sau khi widget unmounted",
        "Đồng bộ dark mode + ngôn ngữ giữa Provider và Firestore",
        "Chuẩn hóa thời gian cho streak (timezone, đầu ngày)",
        "Tối ưu animation trên máy yếu - dùng flutter_animate và Lottie "
        "có thể tốn GPU",
        "Kiểm soát rate limit của Firestore khi nhiều user chơi cùng lúc",
    ], size=13)

    add_text(s, Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.5),
             "Bài học rút ra", size=20, bold=True, color=ACCENT2)
    add_bullets(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(5.0), [
        "Tách trách nhiệm rõ: UI - Provider - Service - Model",
        "Sử dụng async/await chứ không .then() làm code dễ đọc hơn",
        "Stream + StreamBuilder rất mạnh cho realtime UI",
        "Try - catch ở những điểm rủi ro: audio, notification, network",
        "Const widget và keys giúp hiệu năng re-build tốt hơn",
        "Thiết kế design tokens (màu, spacing) giúp UI nhất quán",
    ], size=13)
    add_footer(s, 27, total)


def slide_future(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "26. Hướng phát triển tiếp theo",
                  "Roadmap sau buổi báo cáo")

    items = [
        ("Q3 2026", "Hoàn thiện iOS, ra mắt trên App Store", PRIMARY,
         "Test thiết bị iOS, signing, App Store Connect"),
        ("Q3 2026", "Sign-In Google + Apple", ACCENT,
         "Tăng tỉ lệ đăng ký thành công, giảm ma sát nhập mật khẩu"),
        ("Q4 2026", "Chatbot luyện nói sử dụng Claude / Gemini", ACCENT2,
         "Gọi API LLM hỗ trợ học từ vựng theo ngữ cảnh"),
        ("Q4 2026", "Tính năng xã hội - bằng bạn, thách đấu", GOLD,
         "Realtime Firestore + push notification"),
        ("Q1 2027", "Gói từ vựng động cập nhật từ CMS", PRIMARY_DARK,
         "Backoffice, điều phối nội dung theo cấp độ"),
        ("Q1 2027", "Adaptive learning - SRS", ACCENT,
         "Spaced repetition giúp ghi nhớ lâu"),
    ]
    cw = Inches(6.0)
    ch = Inches(1.65)
    gx = Inches(0.55)
    gy = Inches(1.25)
    for i, (q, name, col, desc) in enumerate(items):
        col_idx = i % 2
        row_idx = i // 2
        x = gx + col_idx * (cw + Inches(0.2))
        y = gy + row_idx * (ch + Inches(0.18))
        card = add_round_rect(s, x, y, cw, ch, WHITE, corner=0.06)
        chip = add_round_rect(s, x + Inches(0.3), y + Inches(0.25),
                              Inches(1.3), Inches(0.42), col, corner=0.5)
        tf = chip.text_frame
        tf.margin_left = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = q
        r.font.bold = True
        r.font.size = Pt(11)
        r.font.color.rgb = WHITE
        r.font.name = FONT_TITLE

        add_text(s, x + Inches(1.75), y + Inches(0.2),
                 cw - Inches(2.0), Inches(0.5),
                 name, size=15, bold=True, color=TEXT_DARK)
        add_text(s, x + Inches(0.3), y + Inches(0.85),
                 cw - Inches(0.5), ch - Inches(0.95),
                 desc, size=12, color=TEXT_MUTED, italic=True)

    add_footer(s, 28, total)


def slide_conclusion(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s, BG_LIGHT)
    add_title_bar(s, "27. Kết luận",
                  "Tổng kết những gì VocabQuest đạt được")

    add_text(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(0.6),
             "VocabQuest - từ ý tưởng tới sản phẩm chạy được",
             size=22, bold=True, color=PRIMARY)

    add_bullets(s, Inches(0.6), Inches(2.0), Inches(12.0), Inches(4.0), [
        "Xây dựng thành công app Flutter đa nền tảng với 14 màn hình, "
        "4 mini-game và hệ thống gamification đầy đủ",
        "Nắm vững và áp dụng lập trình bất đồng bộ (Future, async/await, "
        "Stream) qua các tương tác với Firestore và I/O",
        "Tích hợp Firebase một cách an toàn: Auth - Firestore - Storage, "
        "sử dụng Transaction để bảo toàn toàn vẹn dữ liệu",
        "Tạo trải nghiệm người dùng tốt với animation, âm thanh, "
        "TTS, dark mode, song ngữ, nhắc nhở",
        "Đã chuẩn bị roadmap cụ thể cho các phiên bản tiếp theo",
    ], size=14, line_spacing=1.4)

    chip = add_round_rect(s, Inches(0.6), Inches(6.1),
                          Inches(12.0), Inches(0.7), PRIMARY, corner=0.5)
    tf = chip.text_frame
    tf.margin_left = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Học từ vựng không nhàm chán - VocabQuest biến việc học thành trò chơi"
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = WHITE
    r.font.name = FONT_TITLE
    add_footer(s, 29, total)


def slide_thanks(prs, total, logo_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(s, (0x4C, 0x3F, 0xB8), (0x00, 0xC6, 0xA7), angle=120)

    for cx, cy, cw in [
        (Inches(11.0), Inches(-1.2), Inches(4.5)),
        (Inches(-1.0), Inches(5.5), Inches(4.0)),
        (Inches(8.0), Inches(5.5), Inches(2.0)),
    ]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, cw, cw)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.fill.background()
        sppr = c.fill._xPr
        solid = sppr.find(qn('a:solidFill'))
        if solid is not None:
            srgb = solid.find(qn('a:srgbClr'))
            if srgb is not None:
                a = etree.SubElement(srgb, qn('a:alpha'))
                a.set('val', '14000')

    if logo_path and Path(logo_path).exists():
        try:
            s.shapes.add_picture(str(logo_path),
                                 Inches(5.65), Inches(1.0),
                                 width=Inches(2.0))
        except Exception:
            pass
    add_text(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(1.5),
             "CẢM ƠN THẦY/CÔ\nVÀ CÁC BẠN ĐÃ LẮNG NGHE!",
             font=FONT_TITLE, size=54, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
             "Mong nhận được góp ý để hoàn thiện VocabQuest",
             font=FONT_BODY, size=18, color=WHITE,
             align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
             "Q & A",
             font=FONT_TITLE, size=28, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)


# -------------------- BUILD --------------------
def build(out_path: Path, logo_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    sections = [
        "Giới thiệu đề tài",
        "Lý thuyết Flutter",
        "Lý thuyết Dart",
        "Lập trình bất đồng bộ",
        "Event Loop & Isolate",
        "State management - Provider",
        "Kiến trúc dự án",
        "Thư viện backend",
        "Thư viện UI/Animation",
        "Thư viện khác",
        "Phân tích chức năng (1/2)",
        "Phân tích chức năng (2/2)",
        "Luật tính điểm chung",
        "Tổng quan 4 mini-game",
        "Mini-game Nối từ (Matching)",
        "Mini-game Trắc nghiệm (Quiz)",
        "Mini-game Xếp chữ (Word Puzzle)",
        "Mini-game Lật thẻ (Memory)",
        "Luồng dữ liệu",
        "Sơ đồ dữ liệu",
        "Các màn hình",
        "Bất đồng bộ trong app",
        "Tính năng đã hoàn thành",
        "Hạn chế",
        "Khó khăn & bài học",
        "Hướng phát triển",
        "Kết luận",
    ]
    total = len(sections) + 3

    info = {
        "subject": "Môn học: Lập trình ứng dụng di động - Học kỳ cuối",
        "student": "Vũ Kiều Oanh - 20223610",
        "teacher": "...........................",
        "year": "Năm học 2025 - 2026",
    }

    slide_cover(prs, total, logo_path, info)
    slide_toc(prs, total, sections)
    slide_intro(prs, total)
    slide_flutter_theory(prs, total)
    slide_dart_theory(prs, total)
    slide_async_theory(prs, total)
    slide_event_loop(prs, total)
    slide_state_management(prs, total)
    slide_architecture(prs, total)
    slide_libs_firebase(prs, total)
    slide_libs_ui(prs, total)
    slide_libs_other(prs, total)
    slide_use_case_a(prs, total)
    slide_use_case_b(prs, total)
    slide_scoring_overview(prs, total)
    slide_minigames_overview(prs, total)
    slide_game_matching(prs, total)
    slide_game_quiz(prs, total)
    slide_game_puzzle(prs, total)
    slide_game_memory(prs, total)
    slide_data_flow(prs, total)
    slide_db_schema(prs, total)
    slide_screens_overview(prs, total)
    slide_async_in_app(prs, total)
    slide_done(prs, total)
    slide_undone(prs, total)
    slide_lessons(prs, total)
    slide_future(prs, total)
    slide_conclusion(prs, total)
    slide_thanks(prs, total, logo_path)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print("Saved:", out_path)


if __name__ == "__main__":
    here = Path(__file__).resolve().parent.parent
    logo = here / "assets" / "images" / "backgrounds" / "logo.png"
    out = Path(r"C:/Users/Kieu Anh/Desktop/CD2/VocabQuest_Presentation.pptx")
    build(out, logo)
